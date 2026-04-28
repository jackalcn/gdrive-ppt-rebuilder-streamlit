import html
import io
import json
import os
import re
import subprocess
import time
import urllib.parse
import urllib.request
import zipfile
from datetime import datetime, timedelta, timezone
from pathlib import Path
from tempfile import mkdtemp
from urllib.error import HTTPError

import streamlit as st
from PIL import Image
from pptx import Presentation

PIXEL_TO_EMU = 9525
USER_AGENT = "Mozilla/5.0"
APP_VERSION = "v2026.04.28.1"
APP_REPO = "jackalcn/gdrive-ppt-rebuilder-streamlit"
TW_TZ = timezone(timedelta(hours=8))


def get_release_summary() -> tuple[str, str]:
    try:
        updated_at = datetime.fromtimestamp(Path(__file__).stat().st_mtime, tz=TW_TZ)
        updated_text = updated_at.strftime("%Y-%m-%d %H:%M (UTC+8)")
    except OSError:
        updated_text = "未知"
    return APP_VERSION, updated_text


def get_deploy_commit_short() -> str:
    for key in ("STREAMLIT_GIT_COMMIT", "GITHUB_SHA", "COMMIT_SHA"):
        commit = os.getenv(key, "").strip()
        if commit:
            return commit[:7]

    try:
        result = subprocess.run(
            ["git", "rev-parse", "--short", "HEAD"],
            cwd=Path(__file__).resolve().parent,
            capture_output=True,
            text=True,
            check=True,
            timeout=2,
        )
        commit = result.stdout.strip()
        if commit:
            return commit
    except Exception:  # noqa: BLE001
        pass

    return "未知"


def sanitize_filename(name: str) -> str:
    name = re.sub(r'[<>:"/\\|?*]', "_", name).strip()
    return name or "converted_presentation"


def extract_file_id(url: str) -> str:
    m = re.search(r"/file/d/([a-zA-Z0-9_-]+)", url)
    if not m:
        raise ValueError("無法從連結中解析檔案 ID，請確認是 Google Drive 的 /file/d/... 連結。")
    return m.group(1)


def fetch_text(url: str, timeout: int = 45) -> str:
    req = urllib.request.Request(url, headers={"User-Agent": USER_AGENT})
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        return resp.read().decode("utf-8", "ignore")


def parse_drive_view_metadata(drive_url: str) -> tuple[str, int | None, str, str]:
    file_id = extract_file_id(drive_url)
    view_url = f"https://drive.google.com/file/d/{file_id}/view"

    raw_html = fetch_text(view_url)

    title_match = re.search(r"<title>(.*?)</title>", raw_html, flags=re.IGNORECASE | re.DOTALL)
    if title_match:
        title_raw = html.unescape(title_match.group(1)).strip()
        title = re.sub(r"\s*-\s*Google\s*(雲端硬碟|Drive)\s*$", "", title_raw, flags=re.IGNORECASE)
    else:
        title = file_id

    page_matches = re.findall(r"共\s*(\d+)\s*頁", raw_html)
    total_pages = max(int(x) for x in page_matches) if page_matches else None

    upload_match = re.search(r"viewer/upload\?([^\"']+)", raw_html)
    if not upload_match:
        raise RuntimeError("無法從頁面取得 viewer/upload 參數。")

    upload_path = upload_match.group(0)
    upload_path = upload_path.encode("utf-8").decode("unicode_escape")
    upload_path = upload_path.replace("\\u003d", "=").replace("\\u0026", "&")
    upload_url = urllib.parse.urljoin("https://drive.google.com/", upload_path)

    upload_query = urllib.parse.parse_qs(urllib.parse.urlparse(upload_url).query)
    dsmi = upload_query.get("dsmi", ["texmex"])[0]

    upload_resp = fetch_text(upload_url)
    upload_json_text = upload_resp[5:] if upload_resp.startswith(")]}'") else upload_resp
    try:
        upload_data = json.loads(upload_json_text)
    except json.JSONDecodeError as err:
        raise RuntimeError(f"viewer/upload 回應解析失敗: {err}") from err

    page_path = str(upload_data.get("page", "")).strip()
    if not page_path:
        raise RuntimeError("viewer/upload 回應缺少 page 路徑。")

    page_url = urllib.parse.urljoin("https://drive.google.com/viewer/", page_path)
    page_query = urllib.parse.parse_qs(urllib.parse.urlparse(page_url).query)

    token = page_query.get("id", [""])[0]
    if not token:
        raise RuntimeError("無法從 viewer/page 路徑取得圖像 token。")

    dsmi_from_page = page_query.get("dsmi", [""])[0]
    if dsmi_from_page:
        dsmi = dsmi_from_page

    if total_pages is None:
        meta_path = str(upload_data.get("meta", "")).strip()
        if meta_path:
            meta_url = urllib.parse.urljoin("https://drive.google.com/viewer/", meta_path)
            meta_resp = fetch_text(meta_url)
            meta_json_text = meta_resp[5:] if meta_resp.startswith(")]}'") else meta_resp
            try:
                meta_data = json.loads(meta_json_text)
                pages = meta_data.get("pages")
                if isinstance(pages, int) and pages > 0:
                    total_pages = pages
            except json.JSONDecodeError:
                # Fallback to runtime detection when meta payload is not JSON.
                pass

    return title, total_pages, token, dsmi


def build_viewer_img_url(token: str, dsmi: str, page: int, image_width: int) -> str:
    params = {
        "id": token,
        "dsmi": dsmi,
        "auditContext": "forDisplay",
        "page": str(page),
        "skiphighlight": "true",
        "w": str(image_width),
        "webp": "true",
    }
    return "https://drive.google.com/viewer/img?" + urllib.parse.urlencode(params)


def fetch_slide_image_bytes(
    token: str,
    dsmi: str,
    page: int,
    image_width: int,
    retries: int,
) -> bytes | None:
    url = build_viewer_img_url(token=token, dsmi=dsmi, page=page, image_width=image_width)
    last_err = ""

    for attempt in range(1, retries + 1):
        try:
            req = urllib.request.Request(url, headers={"User-Agent": USER_AGENT})
            with urllib.request.urlopen(req, timeout=45) as resp:
                return resp.read()
        except HTTPError as err:
            if err.code in (400, 404):
                return None
            last_err = f"HTTP {err.code}"
        except Exception as err:  # noqa: BLE001
            last_err = str(err)

        time.sleep(0.6 * attempt)

    raise RuntimeError(f"第 {page + 1} 頁下載失敗: {last_err}")


def detect_total_pages(token: str, dsmi: str, retries: int) -> int:
    probe_width = 64
    first_page = fetch_slide_image_bytes(
        token=token,
        dsmi=dsmi,
        page=0,
        image_width=probe_width,
        retries=retries,
    )
    if first_page is None:
        raise RuntimeError("無法下載第 1 頁，請確認檔案可公開檢視。")

    low = 0
    high = 1
    while True:
        data = fetch_slide_image_bytes(
            token=token,
            dsmi=dsmi,
            page=high,
            image_width=probe_width,
            retries=retries,
        )
        if data is None:
            break
        low = high
        high *= 2
        if high > 4096:
            raise RuntimeError("頁數偵測超出安全上限，請確認來源檔案是否正常。")

    while low + 1 < high:
        mid = (low + high) // 2
        data = fetch_slide_image_bytes(
            token=token,
            dsmi=dsmi,
            page=mid,
            image_width=probe_width,
            retries=retries,
        )
        if data is None:
            high = mid
        else:
            low = mid

    return low + 1


def download_slide_images(
    token: str,
    dsmi: str,
    total_pages: int | None,
    image_width: int,
    retries: int,
    output_dir: Path,
    log,
    update_progress,
) -> tuple[list[Path], tuple[int, int], int]:
    output_dir.mkdir(parents=True, exist_ok=True)
    image_paths: list[Path] = []
    size_set: set[tuple[int, int]] = set()

    if total_pages is None:
        log("未提供總頁數，系統改用 API 自動偵測頁數。")
        total_pages = detect_total_pages(token=token, dsmi=dsmi, retries=retries)
        log(f"頁數偵測完成：共 {total_pages} 頁")

    for idx in range(total_pages):
        data = fetch_slide_image_bytes(
            token=token,
            dsmi=dsmi,
            page=idx,
            image_width=image_width,
            retries=retries,
        )
        if data is None:
            raise RuntimeError(f"第 {idx + 1} 頁回傳空內容，可能為頁數解析異常。")

        image = Image.open(io.BytesIO(data)).convert("RGB")
        size_set.add(image.size)
        out_path = output_dir / f"slide_{idx + 1:03d}.png"
        image.save(out_path, "PNG", optimize=True)
        image_paths.append(out_path)

        update_progress(idx + 1, total_pages)
        if (idx + 1) % 5 == 0 or (idx + 1) == total_pages:
            log(f"已下載 {idx + 1}/{total_pages} 頁")

    if len(size_set) != 1:
        raise RuntimeError(f"下載頁面尺寸不一致: {sorted(size_set)}")

    return image_paths, next(iter(size_set)), total_pages


def build_pptx_from_images(
    image_paths: list[Path],
    size: tuple[int, int],
    output_path: Path,
    log,
    update_progress,
) -> Path:
    width, height = size
    prs = Presentation()
    prs.slide_width = width * PIXEL_TO_EMU
    prs.slide_height = height * PIXEL_TO_EMU
    blank = prs.slide_layouts[6]

    total = len(image_paths)
    for idx, image_path in enumerate(image_paths, 1):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image_path),
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

        update_progress(idx, total)
        if idx % 5 == 0 or idx == total:
            log(f"已建立投影片 {idx}/{total}")

    try:
        prs.save(str(output_path))
        return output_path
    except PermissionError:
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        fallback = output_path.with_stem(output_path.stem + f"_fixed_{ts}")
        prs.save(str(fallback))
        log(f"原檔被占用，已改存為: {fallback.name}")
        return fallback


def build_pdf_from_images(
    image_paths: list[Path],
    output_path: Path,
    log,
    update_progress,
) -> Path:
    pil_images: list[Image.Image] = []
    total = len(image_paths)

    try:
        for idx, image_path in enumerate(image_paths, 1):
            pil_images.append(Image.open(image_path).convert("RGB"))
            update_progress(idx, total)
            if idx % 5 == 0 or idx == total:
                log(f"已準備 PDF 頁面 {idx}/{total}")

        if not pil_images:
            raise RuntimeError("沒有可用影像可輸出 PDF。")

        first = pil_images[0]
        rest = pil_images[1:]
        try:
            first.save(str(output_path), "PDF", save_all=True, append_images=rest)
            return output_path
        except PermissionError:
            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
            fallback = output_path.with_stem(output_path.stem + f"_fixed_{ts}")
            first.save(str(fallback), "PDF", save_all=True, append_images=rest)
            log(f"原檔被占用，已改存為: {fallback.name}")
            return fallback
    finally:
        for img in pil_images:
            img.close()


def build_output_name(naming_mode: str, title: str, custom_name: str, output_format: str) -> str:
    ext = ".pdf" if output_format == "PDF" else ".pptx"
    if naming_mode == "沿用原始檔名":
        return sanitize_filename(title) + ext
    if naming_mode == "原始檔名加日期":
        ts = datetime.now().strftime("%Y%m%d")
        return sanitize_filename(f"{title}_{ts}") + ext
    if not custom_name.strip():
        raise ValueError("請輸入自訂檔名。")
    return sanitize_filename(custom_name.strip()) + ext


def collect_input_links(mode: str, drive_url: str, batch_links_text: str) -> list[str]:
    if mode == "單一連結":
        return [drive_url.strip()] if drive_url.strip() else []
    return [x.strip() for x in batch_links_text.splitlines() if x.strip()]


def history_to_csv_bytes(records: list[dict]) -> bytes:
    headers = [
        "時間",
        "狀態",
        "來源連結",
        "標題",
        "輸出檔名",
        "頁數",
        "尺寸",
        "檔案大小MB",
        "訊息",
    ]
    lines = [",".join(headers)]
    for r in records:
        row = [
            str(r.get("時間", "")),
            str(r.get("狀態", "")),
            str(r.get("來源連結", "")),
            str(r.get("標題", "")),
            str(r.get("輸出檔名", "")),
            str(r.get("頁數", "")),
            str(r.get("尺寸", "")),
            str(r.get("檔案大小MB", "")),
            str(r.get("訊息", "")),
        ]
        escaped = []
        for cell in row:
            if any(ch in cell for ch in [",", '"', "\n"]):
                cell = '"' + cell.replace('"', '""') + '"'
            escaped.append(cell)
        lines.append(",".join(escaped))
    return ("\n".join(lines)).encode("utf-8-sig")


def bundle_zip_bytes(items: list[tuple[str, bytes]]) -> bytes:
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, mode="w", compression=zipfile.ZIP_DEFLATED) as zf:
        for name, data in items:
            zf.writestr(name, data)
    return buf.getvalue()


st.set_page_config(page_title="Drive 簡報轉換器", page_icon="📎", layout="wide")

st.markdown(
    """
<style>
@import url('https://fonts.googleapis.com/css2?family=Noto+Sans+TC:wght@400;500;700&display=swap');
html, body, [class*="css"] { font-family: 'Noto Sans TC', sans-serif; }
.gov-card {
  border: 1px solid #D7DEE3;
  border-radius: 10px;
  padding: 14px;
  background: #FFFFFF;
}
.gov-kpi {
  border-left: 4px solid #0F4C5C;
  padding: 8px 12px;
  background: #F7FAFC;
  border-radius: 6px;
}
.small-note { color: #5F6B73; font-size: 0.9rem; }
</style>
""",
    unsafe_allow_html=True,
)

st.title("Google Drive 簡報轉換平台")
st.caption("專業公務風 | 公開可檢視連結轉換為 PPTX 或 PDF")
app_version, app_updated_at = get_release_summary()
deploy_commit = get_deploy_commit_short()
st.markdown(
    f"<div class='small-note'>版本：{app_version}｜最後更新：{app_updated_at}</div>",
    unsafe_allow_html=True,
)

with st.sidebar:
    st.subheader("任務設定")
    mode = st.radio("作業模式", ["單一連結", "批次連結"], index=0)
    drive_url = ""
    batch_links_text = ""
    if mode == "單一連結":
        drive_url = st.text_input("Google Drive 連結", placeholder="https://drive.google.com/file/d/.../view")
    else:
        batch_links_text = st.text_area(
            "批次連結（每行一個）",
            placeholder="https://drive.google.com/file/d/.../view\nhttps://drive.google.com/file/d/.../view",
            height=150,
        )
    naming_mode = st.radio("輸出檔名規則", ["沿用原始檔名", "原始檔名加日期", "自訂名稱"], index=0)
    custom_name = st.text_input("自訂檔名", value="") if naming_mode == "自訂名稱" else ""
    output_format = st.selectbox("輸出格式", ["PPTX", "PDF"], index=0)
    image_width = st.selectbox("圖片寬度", [1600, 1920, 1280], index=0)
    retries = st.slider("重試次數", 1, 5, 3)
    save_local = st.checkbox("另存至伺服器工作目錄", value=True)
    start = st.button("開始轉換", type="primary", use_container_width=True)
    check_links = st.button("連結健康檢查", use_container_width=True)
    retry_failed = st.button("重試失敗項目", use_container_width=True)
    clear_history = st.button("清除任務歷程", use_container_width=True)
    st.markdown("---")
    st.caption("部署維運")
    st.markdown("推送到 `main` 分支會自動部署，也可按下方連結手動觸發 redeploy。")
    st.link_button(
        "一鍵重新部署（GitHub Actions）",
        f"https://github.com/{APP_REPO}/actions/workflows/streamlit-cloud-redeploy.yml",
        use_container_width=True,
    )

step_col1, step_col2, step_col3, step_col4 = st.columns(4)
step_col1.markdown("<div class='gov-kpi'>1. 解析連結</div>", unsafe_allow_html=True)
step_col2.markdown("<div class='gov-kpi'>2. 下載頁面</div>", unsafe_allow_html=True)
step_col3.markdown("<div class='gov-kpi'>3. 重建檔案</div>", unsafe_allow_html=True)
step_col4.markdown("<div class='gov-kpi'>4. 完成交付</div>", unsafe_allow_html=True)

status_box = st.empty()
download_progress = st.progress(0, text="下載進度: 0%")
build_progress = st.progress(0, text="建檔進度: 0%")
log_box = st.empty()

if "logs" not in st.session_state:
    st.session_state.logs = []
if "history" not in st.session_state:
    st.session_state.history = []
if "link_checks" not in st.session_state:
    st.session_state.link_checks = []

if clear_history:
    st.session_state.history = []
    st.session_state.link_checks = []


def add_log(message: str) -> None:
    ts = datetime.now().strftime("%H:%M:%S")
    st.session_state.logs.append(f"[{ts}] {message}")
    log_box.code("\n".join(st.session_state.logs[-200:]))


run_mode = None
if start:
    run_mode = "start"
elif check_links:
    run_mode = "check"
elif retry_failed:
    run_mode = "retry"

if run_mode:
    st.session_state.logs = []
    log_box.code("")
    download_progress.progress(0, text="下載進度: 0%")
    build_progress.progress(0, text="建檔進度: 0%")

    try:
        if run_mode == "check":
            links = collect_input_links(mode=mode, drive_url=drive_url, batch_links_text=batch_links_text)
            if not links:
                raise ValueError("請至少輸入一個 Google Drive 連結。")

            check_results: list[dict] = []
            ok_count = 0

            for idx, link in enumerate(links, 1):
                status_box.info(f"({idx}/{len(links)}) 系統正在進行連結健康檢查。")
                add_log(f"[{idx}/{len(links)}] 健康檢查開始")

                row = {
                    "時間": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    "狀態": "異常",
                    "來源連結": link,
                    "標題": "",
                    "頁數": "",
                    "訊息": "",
                }

                try:
                    title, total_pages, token, dsmi = parse_drive_view_metadata(link)
                    resolved_pages = total_pages if total_pages is not None else detect_total_pages(token, dsmi, retries)
                    first_page = fetch_slide_image_bytes(
                        token=token,
                        dsmi=dsmi,
                        page=0,
                        image_width=max(256, image_width),
                        retries=retries,
                    )
                    if first_page is None:
                        raise RuntimeError("第一頁下載驗證失敗。")

                    row.update(
                        {
                            "狀態": "可用",
                            "標題": title,
                            "頁數": str(resolved_pages),
                            "訊息": "可正常解析與下載",
                        }
                    )
                    ok_count += 1
                    add_log(f"[{idx}/{len(links)}] 健康檢查通過：{title} / {resolved_pages} 頁")
                except Exception as err:  # noqa: BLE001
                    row["訊息"] = str(err)
                    add_log(f"[{idx}/{len(links)}] 健康檢查失敗：{err}")

                check_results.append(row)

            st.session_state.link_checks = check_results

            if ok_count == len(links):
                status_box.success(f"健康檢查完成：{ok_count}/{len(links)} 可用。")
            elif ok_count > 0:
                status_box.warning(f"健康檢查完成：{ok_count}/{len(links)} 可用，請檢查異常連結。")
            else:
                status_box.error("健康檢查完成：全部連結異常。")

        elif run_mode == "start":
            links = collect_input_links(mode=mode, drive_url=drive_url, batch_links_text=batch_links_text)
            if not links:
                raise ValueError("請至少輸入一個 Google Drive 連結。")

            operation_name = "轉換作業"
        else:
            links = []
            seen = set()
            for rec in st.session_state.history:
                if rec.get("狀態") == "成功":
                    continue
                link = str(rec.get("來源連結", "")).strip()
                if link and link not in seen:
                    seen.add(link)
                    links.append(link)

            if not links:
                raise ValueError("沒有可重試的失敗項目。")

            add_log(f"啟動失敗項目重試，共 {len(links)} 筆")
            operation_name = "重試作業"

        if run_mode != "check":
            artifacts: list[tuple[str, bytes]] = []
            success_count = 0

            for idx, link in enumerate(links, 1):
                row = {
                    "時間": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                    "狀態": "失敗",
                    "來源連結": link,
                    "標題": "",
                    "輸出檔名": "",
                    "頁數": "",
                    "尺寸": "",
                    "檔案大小MB": "",
                    "訊息": "",
                }
                try:
                    status_box.info(f"({idx}/{len(links)}) 系統正在解析連結資訊。")
                    add_log(f"[{idx}/{len(links)}] 開始解析 Drive 連結")
                    title, total_pages, token, dsmi = parse_drive_view_metadata(link)
                    page_text = str(total_pages) if total_pages is not None else "自動偵測"
                    add_log(f"[{idx}/{len(links)}] 解析完成：標題={title}，頁數={page_text}")

                    output_name = build_output_name(naming_mode, title, custom_name, output_format)
                    temp_root = Path(mkdtemp(prefix="drive_ppt_"))
                    image_dir = temp_root / "slides"
                    output_file = temp_root / output_name

                    status_box.info(f"({idx}/{len(links)}) 系統正在下載投影片頁面。")
                    image_paths, size, resolved_pages = download_slide_images(
                        token=token,
                        dsmi=dsmi,
                        total_pages=total_pages,
                        image_width=image_width,
                        retries=retries,
                        output_dir=image_dir,
                        log=add_log,
                        update_progress=lambda done, total: download_progress.progress(
                            int(done * 100 / total), text=f"下載進度: {done}/{total}"
                        ),
                    )

                    status_box.info(f"({idx}/{len(links)}) 系統正在重建 {output_format}。")
                    if output_format == "PDF":
                        final_file = build_pdf_from_images(
                            image_paths=image_paths,
                            output_path=output_file,
                            log=add_log,
                            update_progress=lambda done, total: build_progress.progress(
                                int(done * 100 / total), text=f"建檔進度: {done}/{total}"
                            ),
                        )
                    else:
                        final_file = build_pptx_from_images(
                            image_paths=image_paths,
                            size=size,
                            output_path=output_file,
                            log=add_log,
                            update_progress=lambda done, total: build_progress.progress(
                                int(done * 100 / total), text=f"建檔進度: {done}/{total}"
                            ),
                        )

                    file_bytes = final_file.read_bytes()
                    artifacts.append((final_file.name, file_bytes))
                    success_count += 1

                    row.update(
                        {
                            "狀態": "成功",
                            "標題": title,
                            "輸出檔名": final_file.name,
                            "頁數": str(resolved_pages),
                            "尺寸": f"{size[0]} x {size[1]}",
                            "檔案大小MB": f"{len(file_bytes) / 1024 / 1024:.2f}",
                            "訊息": "完成",
                        }
                    )

                    if save_local:
                        local_path = Path.cwd() / final_file.name
                        try:
                            local_path.write_bytes(file_bytes)
                            add_log(f"已另存至：{local_path}")
                        except PermissionError:
                            ts = datetime.now().strftime("%Y%m%d_%H%M%S")
                            fallback = Path.cwd() / f"{Path(final_file.name).stem}_fixed_{ts}{final_file.suffix}"
                            fallback.write_bytes(file_bytes)
                            add_log(f"原檔被占用，已另存為：{fallback}")

                except Exception as err:  # noqa: BLE001
                    row["訊息"] = str(err)
                    add_log(f"[{idx}/{len(links)}] 錯誤：{err}")

                st.session_state.history.append(row)

            if success_count == len(links):
                status_box.success(f"{operation_name}完成：{success_count}/{len(links)} 成功。")
            elif success_count > 0:
                status_box.warning(f"{operation_name}完成：{success_count}/{len(links)} 成功，請檢查失敗項目。")
            else:
                status_box.error("作業未完成：全部連結皆失敗。")

            if len(artifacts) == 1:
                fname, data = artifacts[0]
                if fname.lower().endswith(".pdf"):
                    label = "下載轉換後 PDF"
                    mime = "application/pdf"
                else:
                    label = "下載轉換後 PPTX"
                    mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                st.download_button(
                    label=label,
                    data=data,
                    file_name=fname,
                    mime=mime,
                    use_container_width=True,
                )
            elif len(artifacts) > 1:
                zip_bytes = bundle_zip_bytes(artifacts)
                zip_name = f"drive_batch_{output_format.lower()}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"
                st.download_button(
                    label="下載批次成果 ZIP",
                    data=zip_bytes,
                    file_name=zip_name,
                    mime="application/zip",
                    use_container_width=True,
                )

    except Exception as err:  # noqa: BLE001
        status_box.error(f"作業未完成：{err}")
        add_log(f"錯誤：{err}")

if st.session_state.history:
    st.subheader("任務歷程")
    st.dataframe(st.session_state.history, use_container_width=True)
    csv_bytes = history_to_csv_bytes(st.session_state.history)
    st.download_button(
        label="下載任務歷程 CSV",
        data=csv_bytes,
        file_name=f"task_history_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv",
        mime="text/csv",
        use_container_width=True,
    )

if st.session_state.link_checks:
    st.subheader("連結健康檢查結果")
    st.dataframe(st.session_state.link_checks, use_container_width=True)

with st.expander("作業說明", expanded=False):
    st.markdown(
        """
- 僅支援 Google Drive 可公開檢視連結。
- 若原始檔限制下載，本工具會透過檢視頁圖像重建為新 PPTX。
- 若同名檔案被占用，系統會自動附加 `_fixed_時間戳` 後存檔。
"""
    )

st.markdown("---")
st.caption(f"部署提交：{deploy_commit}")
